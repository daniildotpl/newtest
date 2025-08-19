import pandas as pd
from ishch.connector import my_connection
from pptx import Presentation



LIST_OF_TAGS = [
    'reporting_year', 
    'prevent_year', 
    'reporting_region', 
    'region_2021', 
    'kids_2021', 
    'adults_2021', 
    'count_org_2021', 
    'count_org_rank_2021', 
    'doctors_2021', 
    'doctors_rank_2021', 
    'nurses_2021', 
    'nurses_rank_2021', 
    'region_2022', 
    'kids_2022', 
    'adults_2022', 
    'count_org_2022', 
    'count_org_rank_2022', 
    'doctors_2022', 
    'doctors_rank_2022', 
    'nurses_2022', 
    'nurses_rank_2022'
]


def get_payload(engine, reporting_region):

    query = f"""
        WITH preset AS (
            SELECT 
                year
                , region
                , SUM(CASE 
                    WHEN tip_mo = 'Детские поликлиники' AND name_value = 'Число организаций, ед'
                    THEN TRY_CAST(REPLACE(REPLACE(value, ' ', ''), ',', '.') AS DECIMAL(15,0))
                    ELSE 0 END) AS kids
                , SUM(CASE
                    WHEN tip_mo = 'Поликлиники' AND name_value = 'Число организаций, ед'
                    THEN TRY_CAST(REPLACE(REPLACE(value, ' ', ''), ',', '.') AS DECIMAL(15,0))
                    ELSE 0 END) AS adults
                , SUM(CASE 
                    WHEN name_value = 'Число организаций, ед'
                    THEN TRY_CAST(REPLACE(REPLACE(value, ' ', ''), ',', '.') AS DECIMAL(15,0)) 
                    ELSE 0 END) AS count_org
                , SUM(CASE 
                    WHEN name_value = 'Число должностей врачей, ед занятых' 
                    THEN TRY_CAST(REPLACE(REPLACE(value, ' ', ''), ',', '.') AS DECIMAL(15,1)) 
                    ELSE 0 END) AS doctors
                , SUM(CASE 
                    WHEN name_value = 'Число должностей среднего медперсонала, ед занятых' 
                    THEN TRY_CAST(REPLACE(REPLACE(value, ' ', ''), ',', '.') AS DECIMAL(15,1)) 
                    ELSE 0 END) AS nurses
            FROM statinfo
            GROUP BY year, region        
        ),
        ranked AS (
            SELECT 
                year
                , region
                , kids
                , adults
                , count_org
                , doctors
                , nurses
                , DENSE_RANK() OVER (ORDER BY count_org DESC) as count_org_rank
                , DENSE_RANK() OVER (ORDER BY doctors DESC) as doctors_rank
                , DENSE_RANK() OVER (ORDER BY nurses DESC) as nurses_rank
            FROM preset
        )
        SELECT * 
        FROM ranked
        WHERE year IN (2021, 2022) 
        AND region = {reporting_region}
        ;
    """
    payload = pd.read_sql(query, engine)
    print(payload)
    return payload


def get_context(payload):
    context = {}
    context['reporting_year'] = int(payload.loc[1]['year'])
    context['prevent_year'] = int(payload.loc[1]['year']) - 1
    context['reporting_region'] = int(payload.loc[1]['region'])
    # ---
    payload_as_dict = payload.set_index('year').to_dict('index')
    for year, dictt in payload_as_dict.items():
        for key, value in dictt.items():
            key = key + '_' + str(year)
            print(key)
            context[key] = value
    # ---
    context['kids'] = payload.loc[1]['kids']
    context['adults'] = payload.loc[1]['adults']
    print(f'context: {context}')
    return context


def make_powerp(context):

    prs = Presentation('ishch/template.pptx')  

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # print()
                # print('has_text_frame')
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # print()
                        # print(f'run: {run.text}')
                        for tag in LIST_OF_TAGS:
                            if tag in run.text:
                                # print(f'find: {tag}')
                                try:
                                    value = context[tag]
                                    # print(f'context value: {value}')        
                                    run.text = run.text.replace(tag, str(context[tag]))
                                    # print('success')
                                except Exception as error:
                                    print(f'Ошибка в Run: {error}')
            if shape.has_table:
                # print()
                # print('has_table')
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        # print()
                        # print(f'cell: {cell.text}')
                        for tag in LIST_OF_TAGS:
                            if tag in cell.text:
                                print(f'find: {tag}')
                                try:
                                    value = context[tag]
                                    # print(f'context value: {value}')        
                                    cell.text = cell.text.replace(tag, str(context[tag]))
                                    # print('success')
                                except Exception as error:
                                    print(f'Ошибка в Cell: {error}')

    prs.save('generated_presentation.pptx')


if __name__ == "__main__":
    payload = get_payload(my_connection(), 2336)
    context = get_context(payload)
    make_powerp(context)




